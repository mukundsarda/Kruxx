from flask import Flask, render_template, request, redirect, url_for, session, jsonify
from flask_cors import CORS
from transformers import pipeline
import os
import json
import ast
import textract
import requests
from ocr_processing import process_image
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import TextFormatter
from workers import txt2questions
from pptx import Presentation
from dotenv import load_dotenv
from googletrans import Translator
from bs4 import BeautifulSoup
from youtubesearchpython import VideosSearch
from yt_dlp import YoutubeDL
import re
import urllib.parse
import urllib.request
import numpy as np
import google.generativeai as genai
from textblob import TextBlob
from googleapiclient.discovery import build
from sklearn.metrics.pairwise import cosine_similarity
from sklearn.feature_extraction.text import TfidfVectorizer
app = Flask(__name__)
CORS(app)
load_dotenv()

app.secret_key = os.getenv('secret_key')
YOUTUBE_API_KEY = os.getenv('YOUTUBE_API_KEY')
gemini_api_key=os.getenv('gemini_api_key')
genai.configure(api_key=gemini_api_key)
model = genai.GenerativeModel(model_name="gemini-1.5-pro-latest")
# print(app.secret_key)

bart_summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
bert_summarizer = pipeline("summarization", model="bert-base-uncased")


@app.route('/')
def main():
    return render_template("index.html")


@app.route('/upload-image', methods=['POST'])
def upload_image():
    # Retrieve the uploaded file
    if 'image' in request.files and request.files['image'].filename != '':
        f = request.files['image']
        filename = f.filename
        f.save(filename)

        if filename.lower().endswith(('.jpg', '.jpeg', '.png')):

            # Retrieve radio button selections
            performance = request.form.get('performance')  # Faster or Better
            summary_type = request.form.get('summaryType')
            summary_length = request.form.get('summaryLength')

            output_text = process_image(filename, performance)
            output_filename = f"{filename}_output.txt"

            with open(output_filename, 'w') as output_file:
                output_file.write(output_text)

            with open(output_filename, 'r', encoding="utf-8") as file:
                text = file.read()

            max_input_length = 512  # Model's input token limit

            if summary_length == 'short':
                max_output_length = 50
            elif summary_length == 'medium':
                max_output_length = 130
            elif summary_length == 'long':
                max_output_length = 250
            else:
                max_output_length = 130

            truncated_text = text[:max_input_length]

            try:
                if summary_type == 'abstractive':
                    summary = bart_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                elif summary_type == 'extractive':
                    summary = bert_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                else:
                    summary = "Invalid summary type selected."

                summary_filename = f"{filename}_summary.txt"
                
                with open(summary_filename, 'w', encoding="utf-8") as summary_file:
                    summary_file.write(summary)

                message = f"Summary generated and saved to {summary_filename} using {summary_type.upper()} with {summary_length} length."
                return jsonify(success=True, message=message, summary=summary)  # Send summary to frontend
            except Exception as e:
                return jsonify(success=False, message=f"Error during summarization: {e}")

        else:
            return jsonify(success=False, message="File uploaded successfully but it is not an image.")
        
    return jsonify(success=False, message="No file was uploaded.")



@app.route('/upload-pdf', methods=['POST'])
def upload_pdf():
    if 'pdf' in request.files and request.files['pdf'].filename != '':
        pdf_file = request.files['pdf']
        filename = pdf_file.filename
        pdf_file.save(filename)
        
        if filename.lower().endswith(('pdf')):
            extracted_text = textract.process(filename).decode('utf-8')
            text_filename = f"{filename}_text.txt"

            with open(text_filename, "w", encoding="utf-8") as text_file:
                text_file.write(extracted_text)
    
            summary_type = request.form.get('summaryType')
            summary_length = request.form.get('summaryLength')

            # with open(text_filename, 'w') as output_file:
            #     output_file.write(text)

            # with open(text_filename, 'r', encoding="utf-8") as file:
            #     text = file.read()

            max_input_length = 512  # Model's input token limit

            if summary_length == 'short':
                max_output_length = 50
            elif summary_length == 'medium':
                max_output_length = 130
            elif summary_length == 'long':
                max_output_length = 250
            else:
                max_output_length = 130

            truncated_text = extracted_text[:max_input_length]

            try:
                if summary_type == 'abstractive':
                    summary = bart_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                elif summary_type == 'extractive':
                    summary = bert_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                else:
                    summary = "Invalid summary type selected."

                summary_filename = f"{filename}_summary.txt"
                
                with open(summary_filename, 'w', encoding="utf-8") as summary_file:
                    summary_file.write(summary)

                message = f"Summary generated and saved to {summary_filename} using {summary_type.upper()} with {summary_length} length."
                return jsonify(success=True, message=message, summary=summary)  # Send summary to frontend
            except Exception as e:
                return jsonify(success=False, message=f"Error during summarization: {e}")
            
        else:
            return jsonify(success=False, message="File uploaded successfully but it is not a PDF.")
        
    return jsonify(success=False, message="No file was uploaded.")



@app.route('/upload-ppt', methods=['POST'])
def upload_ppt():
    if 'ppt' in request.files and request.files['ppt'].filename != '':
        ppt_file = request.files['ppt']
        filename = ppt_file.filename
        ppt_file.save(filename)

        if filename.lower().endswith(('pptx')):

            def extract_text_from_ppt(filename):
                prs = Presentation(filename)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                text.append(paragraph.text)
                return "\n".join(text)

            extracted_text = extract_text_from_ppt(filename)
            text_filename = f"{filename}_text.txt"

            print(text_filename, "\n\n\n")

            with open(text_filename, "w", encoding="utf-8") as text_file:
                text_file.write(extracted_text)
    

            summary_type = request.form.get('summaryType')
            summary_length = request.form.get('summaryLength')
            
            max_input_length = 512  # Model's input token limit

            if summary_length == 'short':
                max_output_length = 50
            elif summary_length == 'medium':
                max_output_length = 130
            elif summary_length == 'long':
                max_output_length = 250
            else:
                max_output_length = 130

            truncated_text = extracted_text[:max_input_length]

            try:
                if summary_type == 'abstractive':
                    summary = bart_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                elif summary_type == 'extractive':
                    summary = bert_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                else:
                    summary = "Invalid summary type selected."

                summary_filename = f"{filename}_summary.txt"
                
                with open(summary_filename, 'w', encoding="utf-8") as summary_file:
                    summary_file.write(summary)

                message = f"Summary generated and saved to {summary_filename} using {summary_type.upper()} with {summary_length} length."
                return jsonify(success=True, message=message, summary=summary)  # Send summary to frontend
            except Exception as e:
                return jsonify(success=False, message=f"Error during summarization: {e}")
            
        else:
            return jsonify(success=False, message="File uploaded successfully but it is not a PPT.")
        
    return jsonify(success=False, message="No file was uploaded.")



@app.route('/upload-doc', methods=['POST'])
def upload_doc():
    if 'doc' in request.files and request.files['doc'].filename != '':
        doc_file = request.files['doc']
        filename = doc_file.filename
        doc_file.save(filename)
        
        if filename.lower().endswith(('doc', 'docx')):
            extracted_text = textract.process(filename).decode('utf-8')
            text_filename = f"{filename}_text.txt"
            
            with open(text_filename, "w", encoding="utf-8") as text_file:
                text_file.write(extracted_text)

            summary_type = request.form.get('summaryType')
            summary_length = request.form.get('summaryLength')
            
            max_input_length = 512  # Model's input token limit

            if summary_length == 'short':
                max_output_length = 50
            elif summary_length == 'medium':
                max_output_length = 130
            elif summary_length == 'long':
                max_output_length = 250
            else:
                max_output_length = 130

            truncated_text = extracted_text[:max_input_length]

            try:
                if summary_type == 'abstractive':
                    summary = bart_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                elif summary_type == 'extractive':
                    summary = bert_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                else:
                    summary = "Invalid summary type selected."

                summary_filename = f"{filename}_summary.txt"
                
                with open(summary_filename, 'w', encoding="utf-8") as summary_file:
                    summary_file.write(summary)

                message = f"Summary generated and saved to {summary_filename} using {summary_type.upper()} with {summary_length} length."
                return jsonify(success=True, message=message, summary=summary)  # Send summary to frontend
            except Exception as e:
                return jsonify(success=False, message=f"Error during summarization: {e}")
            
        else:
            return jsonify(success=False, message="File uploaded successfully but it is not a DOC.")
        
    return jsonify(success=False, message="No file was uploaded.")



# @app.route('/quiz-image', methods=['POST'])
# def quiz_image():



@app.route('/quiz-pdf', methods=['POST'])
def quiz_pdf():
    try:
        if 'pdf' in request.files and request.files['pdf'].filename != '':
            pdf_file = request.files['pdf']
            filename = pdf_file.filename
            pdf_file.save(filename)

            if filename.lower().endswith('.pdf'):
                # Extract text from the PDF
                extracted_text = textract.process(filename).decode('utf-8')
                text_filename = f"{filename}_text.txt"

                # Save the extracted text to a .txt file
                with open(text_filename, "w", encoding="utf-8") as text_file:
                    text_file.write(extracted_text)

                # Generate quiz questions from the text file
                with open(text_filename, "r", encoding="utf-8") as file:
                    text = file.read()

                global questions
                
                prompt = f"""
                Generate exactly 10 multiple-choice questions (MCQs) from the text below.
                Return only the dictionary itself, WITHOUT ASSIGNING IT TO A VARIABLE. No code fences, no explanation.
                Return the result in the following Python dictionary format:

                {{
                    1: {{
                        "question": "...",
                        "answer": "...",
                        "options": ["...", "...", "...", "..."]
                    }},
                    2: {{ ... }},
                    ...
                }}

                Text:
                {text}
                """

                # Generate MCQs using Gemini
                response = model.generate_content(prompt)
                content = response.text.strip()

                # Debug: Log what we got from Gemini
                print("RAW GEMINI RESPONSE:")
                print(content)

                if not content:
                    return jsonify({'success': False, 'message': 'Gemini returned an empty response.'}), 500

                # Remove markdown code fences (```json or ```)
                content_clean = re.sub(r'^```[a-zA-Z]*|```$', '', content).strip()

                # Extract dictionary block only
                match = re.search(r'\{.*\}', content_clean, re.DOTALL)
                if match:
                    content_clean = match.group(0)
                else:
                    return jsonify({'success': False, 'message': 'Could not find dictionary content in Gemini response.'}), 500

                try:
                    questions = ast.literal_eval(content_clean)
                except Exception as parse_err:
                    return jsonify({
                        'success': False,
                        'message': f'Failed to parse dictionary from Gemini response: {parse_err}'
                    }), 500


                if questions:
                    return jsonify({
                        "success": True,
                        "message": "Quiz generated successfully!",
                        "questions": questions,
                    })
                else:
                    return jsonify({
                        "success": False,
                        "message": "Failed to generate quiz from the provided text.",
                    })
            else:
                return jsonify({
                    "success": False,
                    "message": "Uploaded file is not a PDF.",
                })
        else:
            return jsonify({
                "success": False,
                "message": "No file was uploaded.",
            })
    except Exception as e:
        return jsonify({
            "success": False,
            "message": f"Error processing quiz from PDF: {e}",
        })



@app.route('/quiz-doc', methods=['POST'])
def quiz_doc():
    try:
        if 'doc' in request.files and request.files['doc'].filename != '':
            doc_file = request.files['doc']
            filename = doc_file.filename
            doc_file.save(filename)

            if filename.lower().endswith(('.doc','.docx')):
                # Extract text from the PDF
                extracted_text = textract.process(filename).decode('utf-8')
                text_filename = f"{filename}_text.txt"

                # Save the extracted text to a .txt file
                with open(text_filename, "w", encoding="utf-8") as text_file:
                    text_file.write(extracted_text)

                # Generate quiz questions from the text file
                with open(text_filename, "r", encoding="utf-8") as file:
                    text = file.read()

                global questions
                prompt = f"""
                Generate exactly 10 multiple-choice questions (MCQs) from the text below.
                Return only the dictionary itself, WITHOUT ASSIGNING IT TO A VARIABLE. No code fences, no explanation.
                Return the result in the following Python dictionary format:

                {{
                    1: {{
                        "question": "...",
                        "answer": "...",
                        "options": ["...", "...", "...", "..."]
                    }},
                    2: {{ ... }},
                    ...
                }}

                Text:
                {text}
                """

                # Generate MCQs using Gemini
                response = model.generate_content(prompt)
                content = response.text.strip()

                # Debug: Log what we got from Gemini
                print("RAW GEMINI RESPONSE:")
                print(content)

                if not content:
                    return jsonify({'success': False, 'message': 'Gemini returned an empty response.'}), 500

                # Remove markdown code fences (```json or ```)
                content_clean = re.sub(r'^```[a-zA-Z]*|```$', '', content).strip()

                # Extract dictionary block only
                match = re.search(r'\{.*\}', content_clean, re.DOTALL)
                if match:
                    content_clean = match.group(0)
                else:
                    return jsonify({'success': False, 'message': 'Could not find dictionary content in Gemini response.'}), 500

                try:
                    questions = ast.literal_eval(content_clean)
                except Exception as parse_err:
                    return jsonify({
                        'success': False,
                        'message': f'Failed to parse dictionary from Gemini response: {parse_err}'
                    }), 500

                if questions:
                    return jsonify({
                        "success": True,
                        "message": "Quiz generated successfully!",
                        "questions": questions,
                    })
                else:
                    return jsonify({
                        "success": False,
                        "message": "Failed to generate quiz from the provided text.",
                    })
            else:
                return jsonify({
                    "success": False,
                    "message": "Uploaded file is not a DOC.",
                })
        else:
            return jsonify({
                "success": False,
                "message": "No file was uploaded.",
            })
    except Exception as e:
        return jsonify({
            "success": False,
            "message": f"Error processing quiz from DOC: {e}",
        })



@app.route('/quiz-ppt', methods=['POST'])
def quiz_ppt():
    try:
        if 'ppt' in request.files and request.files['ppt'].filename != '':
            ppt_file = request.files['ppt']
            filename = ppt_file.filename
            ppt_file.save(filename)

            if filename.lower().endswith('.pptx'):
                # Extract text from the PDF
                def extract_text_from_ppt(filename):
                    prs = Presentation(filename)
                    text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    text.append(paragraph.text)
                    return "\n".join(text)

                extracted_text = extract_text_from_ppt(filename)
                text_filename = f"{filename}_text.txt"

                print(text_filename, "\n\n\n")

                with open(text_filename, "w", encoding="utf-8") as text_file:
                    text_file.write(extracted_text)
                
                # Generate quiz questions from the text file
                with open(text_filename, "r", encoding="utf-8") as file:
                    text = file.read()
                global questions
                prompt = f"""
                Generate exactly 10 multiple-choice questions (MCQs) from the text below.
                Return only the dictionary itself, WITHOUT ASSIGNING IT TO A VARIABLE. No code fences, no explanation.
                Return the result in the following Python dictionary format:

                {{
                    1: {{
                        "question": "...",
                        "answer": "...",
                        "options": ["...", "...", "...", "..."]
                    }},
                    2: {{ ... }},
                    ...
                }}

                Text:
                {text}
                """

                # Generate MCQs using Gemini
                response = model.generate_content(prompt)
                content = response.text.strip()

                # Debug: Log what we got from Gemini
                print("RAW GEMINI RESPONSE:")
                print(content)

                if not content:
                    return jsonify({'success': False, 'message': 'Gemini returned an empty response.'}), 500

                # Remove markdown code fences (```json or ```)
                content_clean = re.sub(r'^```[a-zA-Z]*|```$', '', content).strip()

                # Extract dictionary block only
                match = re.search(r'\{.*\}', content_clean, re.DOTALL)
                if match:
                    content_clean = match.group(0)
                else:
                    return jsonify({'success': False, 'message': 'Could not find dictionary content in Gemini response.'}), 500

                try:
                    questions = ast.literal_eval(content_clean)
                except Exception as parse_err:
                    return jsonify({
                        'success': False,
                        'message': f'Failed to parse dictionary from Gemini response: {parse_err}'
                    }), 500

                if questions:
                    return jsonify({
                        "success": True,
                        "message": "Quiz generated successfully!",
                        "questions": questions,
                    })
                else:
                    return jsonify({
                        "success": False,
                        "message": "Failed to generate quiz from the provided text.",
                    })
            else:
                return jsonify({
                    "success": False,
                    "message": "Uploaded file is not a PPT.",
                })
        else:
            return jsonify({
                "success": False,
                "message": "No file was uploaded.",
            })
    except Exception as e:
        return jsonify({
            "success": False,
            "message": f"Error processing quiz from PPT: {e}",
        })




@app.route('/process_image', methods=['POST'])
def process_image_route():
    filename = request.form.get('filename')
    ocr_option = request.form.get('ocr_option')
    
    output_text = process_image(filename, ocr_option)
    output_filename = f"{filename}_output.txt"
    with open(output_filename, 'w') as output_file:
        output_file.write(output_text)
    
    return render_template("select_summary.html", filename=output_filename)



@app.route('/quizup', methods=['GET', 'POST'])
def quizup():
    try:
        global questions
        questions = {}

        if 'quiz' not in request.files or request.files['quiz'].filename == '':
            return jsonify({'success': False, 'message': 'No file was uploaded.'}), 400

        quiz_file = request.files['quiz']
        filename = quiz_file.filename
        quiz_file.save(filename)

        if not filename.lower().endswith('.txt'):
            return jsonify({'success': False, 'message': 'Uploaded file is not a .txt file.'}), 400

        with open(filename, 'r', encoding='utf-8') as file:
            text = file.read()

        prompt = f"""
        Generate exactly 10 multiple-choice questions (MCQs) from the text below.
        Return only the dictionary itself, WITHOUT ASSIGNING IT TO A VARIABLE. No code fences, no explanation.
        Return the result in the following Python dictionary format:

        {{
            1: {{
                "question": "...",
                "answer": "...",
                "options": ["...", "...", "...", "..."]
            }},
            2: {{ ... }},
            ...
        }}

        Text:
        {text}
        """

        # Generate MCQs using Gemini
        response = model.generate_content(prompt)
        content = response.text.strip()

        # Debug: Log what we got from Gemini
        print("RAW GEMINI RESPONSE:")
        print(content)

        if not content:
            return jsonify({'success': False, 'message': 'Gemini returned an empty response.'}), 500

        # Remove markdown code fences (```json or ```)
        content_clean = re.sub(r'^```[a-zA-Z]*|```$', '', content).strip()

        # Extract dictionary block only
        match = re.search(r'\{.*\}', content_clean, re.DOTALL)
        if match:
            content_clean = match.group(0)
        else:
            return jsonify({'success': False, 'message': 'Could not find dictionary content in Gemini response.'}), 500

        try:
            questions = ast.literal_eval(content_clean)
        except Exception as parse_err:
            return jsonify({
                'success': False,
                'message': f'Failed to parse dictionary from Gemini response: {parse_err}'
            }), 500

        return jsonify({
            "success": True,
            "message": "Quiz generated successfully!"
            # "questions": questions   # Uncomment this if you want to return the questions too
        })

    except Exception as e:
        return jsonify({'success': False, 'message': f'Unexpected error: {e}'}), 500



@app.route('/success', methods=['POST'])
def success():
    if request.method == 'POST':
        option = request.form.get('option')

        # if option == 'quiz':
        #     return render_template("select_quiz.html", filename=request.form.get('filename'))
        
        if option == 'image' and 'image' in request.files and request.files['image'].filename != '':
            f = request.files['image']
            filename = f.filename
            f.save(filename)
            
            if filename.lower().endswith(('.jpg', '.jpeg', '.png')):
                return render_template("select_ocr.html", filename=filename)
            else:
                return render_template("ack.html", message="File uploaded successfully but it is not an image.")
        
        # Handle URL link submission
        elif option == 'link' and request.form.get('link'):
            link = request.form['link']
            try:
                # Fetch HTML content of the URL
                response = requests.get(link)
                response.raise_for_status()
                
                # Save HTML content to a file
                html_filename = "downloaded_link.html"
                with open(html_filename, "w", encoding="utf-8") as html_file:
                    html_file.write(response.text)
                
                # Process the saved HTML file with Textract
                extracted_text = textract.process(html_filename, method='html').decode('utf-8')
                text_filename = "link_extracted_text.txt"
                with open(text_filename, "w", encoding="utf-8") as text_file:
                    text_file.write(extracted_text)
                return render_template("select_summary.html", filename=text_filename)
            except requests.exceptions.RequestException as e:
                message = f"Error fetching the URL: {e}"
            except textract.exceptions.ShellError as e:
                message = f"Error processing the HTML file with Textract: {e}"
            except Exception as e:
                message = f"An unexpected error occurred: {e}"
            return render_template("ack.html", message=message)

        # Handle PDF file upload
        elif option == 'pdf' and 'pdf' in request.files and request.files['pdf'].filename != '':
            pdf_file = request.files['pdf']
            filename = pdf_file.filename
            pdf_file.save(filename)
            if filename.lower().endswith(('pdf')):
                try:
                    extracted_text = textract.process(filename).decode('utf-8')
                    text_filename = f"{filename}_text.txt"
                    with open(text_filename, "w", encoding="utf-8") as text_file:
                        text_file.write(extracted_text)
                    return render_template("select_summary.html", filename=text_filename)
                except Exception as e:
                    message = f"Error processing the PDF file: {e}"
                    return render_template("ack.html", message=message)
            else:
                return render_template("ack.html", message="File uploaded successfully but it is not a pdf.")
        
        # Handle DOC file upload
        elif option == 'doc' and 'doc' in request.files and request.files['doc'].filename != '':
            doc_file = request.files['doc']
            filename = doc_file.filename
            doc_file.save(filename)
            if filename.lower().endswith(('doc', 'docx')):
                try:
                    extracted_text = textract.process(filename).decode('utf-8')
                    text_filename = f"{filename}_text.txt"
                    with open(text_filename, "w", encoding="utf-8") as text_file:
                        text_file.write(extracted_text)
                    return render_template("select_summary.html", filename=text_filename)
                except Exception as e:
                    message = f"Error processing the DOC file: {e}"
                    return render_template("ack.html", message=message)
            else:
                return render_template("ack.html", message="File uploaded successfully but it is not a doc.")

        
        elif option == 'quiz' and 'quiz' in request.files and request.files['quiz'].filename != '':

            # try:
            #     global questions
            #     questions = {}

            #     quiz_file = request.files['quiz']
            #     filename = quiz_file.filename
            #     quiz_file.save(filename)

            #     if filename.lower().endswith('.txt'):
            #         with open(filename, 'r') as file:
            #             text = file.read()
            #         questions = txt2questions(text)

            #         if text:
            #             return jsonify({
            #                 "success": True,
            #                 "message": "Quiz generated successfully!"
            #             })
            #     else:
            #         return jsonify({
            #             "success": False,
            #             "message": "Invalid file type. Please upload a text file."
            #         })
            # except Exception as e:
            #     return jsonify({
            #         "success": False,
            #         "message": f"Error generating quiz: {e}"
            #     })






            UPLOAD_STATUS = False
            global questions
            questions = dict()

            quiz_file = request.files['quiz']
            # # print("\n\n\n\n\n")
            # # print("Quiz File?:", quiz_file)
            # # print("\n\n\n\n\n")
            filename = quiz_file.filename
            # # print("\n\n\n\n\n")
            # # print("Filename:", filename)
            # # print("\n\n\n\n\n")
            quiz_file.save(filename)

            if filename.lower().endswith(('txt')):
                try:
                    with open(filename, 'r') as file:
                        text = file.read()

                    questions = txt2questions(text)
                    
                    # print("\n\n\n\n\n")
                    # print(questions)
                    # print("\n\n\n\n\n")

                # File upload + convert success
                    if text is not None:
                        UPLOAD_STATUS = True

                    # print("\n\n\n\n\n")
                    # print("Reached Debug Point 1")
                    # print("\n\n\n\n\n")

                    # return redirect("/quiz")
                    return render_template('quiz.html', uploaded=UPLOAD_STATUS, questions=questions, size=len(questions))
                except Exception as e:
                    message = f"Error processing the TXT file: {e}"
                    return render_template("ack.html", message=message)
            else:
                return render_template("ack.html", message="File uploaded successfully but it is not a text file.")
        
        # Handle YouTube video URL
        elif option == 'video' and request.form.get('video'):
            video_url = request.form['video']
            video_id = video_url.split('v=')[-1]
            try:
                transcript = YouTubeTranscriptApi.get_transcript(video_id)
                formatter = TextFormatter()
                formatted_transcript = formatter.format_transcript(transcript)
                
                transcript_filename = f"{video_id}_transcript.txt"
                with open(transcript_filename, 'w', encoding="utf-8") as transcript_file:
                    transcript_file.write(formatted_transcript)
                
                return render_template("select_summary.html", filename=transcript_filename)
            except Exception as e:
                message = f"Error fetching transcript for the video: {e}"
                return render_template("ack.html", message=message)
            
        elif option == 'ppt' and 'ppt' in request.files and request.files['ppt'].filename != '':
            ppt_file = request.files['ppt']
            filename = ppt_file.filename
            ppt_file.save(filename)

            if filename.lower().endswith(('pptx')):

                def extract_text_from_ppt(filename):
                    prs = Presentation(filename)
                    text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    text.append(paragraph.text)
                    return "\n".join(text)

                try:
                    extracted_text = extract_text_from_ppt(filename)
                    text_filename = f"{filename}_text.txt"
                    print(text_filename, "\n\n\n")
                    with open(text_filename, "w", encoding="utf-8") as text_file:
                        text_file.write(extracted_text)
                    return render_template("select_summary.html", filename=text_filename)
                except Exception as e:
                    message = f"Error processing the PPT file: {e}"
                    return render_template("ack.html", message=message)
                # except Exception as e:
                #     message = f"Error processing the PPT file: {e}"

                
                # print("\n\n\n", filename, "\n\n\n")
                
                # try:
                #     extracted_text = textract.process(filename).decode('utf-8')
                #     text_filename = f"{filename}_text.txt"
                #     print(text_filename, "\n\n\n")
                #     with open(text_filename, "w", encoding="utf-8") as text_file:
                #         text_file.write(extracted_text)
                #     return render_template("select_summary.html", filename=text_filename)
                # except Exception as e:
                #     message = f"Error processing the PPT file: {e}"
                #     return render_template("ack.html", message=message)
            
            else:
                return render_template("ack.html", message="File uploaded successfully but it is not in the form of pptx. Please make sure it is a .pptx file")

        else:
            message = "Please select a valid option and submit the required information."
            return render_template("ack.html", message=message)


def extract_readable_text(html_content):
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Remove script and style elements
    for script in soup(["script", "style", "header", "footer", "nav"]):
        script.decompose()
    
    # Get text and remove extra whitespace
    text = soup.get_text(separator=' ')
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    text = ' '.join(lines)
    
    # Remove multiple spaces and special characters
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^\w\s.,!?-]', '', text)
    
    return text.strip()

@app.route('/link-upload', methods=['POST'])
def link_upload():
    try:
        # Parse incoming JSON data
        data = request.get_json()
        if not data:
            return jsonify({"success": False, "message": "No data received"}), 400
            
        link = data.get('link')
        summary_type = data.get('summary_type')
        summary_length = data.get('summary_length')

        # Validate input
        if not link:
            return jsonify({"success": False, "message": "Link is required."}), 400
        if not summary_type:
            summary_type = 'abstractive'  # default value
        if not summary_length:
            summary_length = 'medium'  # default value

        try:
            # Make request to get the webpage content
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            response = requests.get(link, headers=headers, timeout=10)
            response.raise_for_status()  # Raise an exception for bad status codes
        except requests.RequestException as e:
            return jsonify({"success": False, "message": f"Error fetching website: {str(e)}"}), 400

        # Extract readable content from HTML
        try:
            extracted_text = extract_readable_text(response.text)
            if not extracted_text:
                return jsonify({"success": False, "message": "Could not extract readable content from the webpage"}), 400
        except Exception as e:
            return jsonify({"success": False, "message": f"Error extracting text: {str(e)}"}), 400

        text_filename = "link_extracted_text.txt"
        with open(text_filename, "w", encoding="utf-8") as text_file:
            text_file.write(extracted_text)

        max_input_length = 512

        if summary_length == 'short':
            max_output_length = 50
        elif summary_length == 'medium':
            max_output_length = 130
        elif summary_length == 'long':
            max_output_length = 250
        else:
            max_output_length = 130

        truncated_text = extracted_text[:max_input_length]

        try:
            if summary_type == 'abstractive':
                summary = bart_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
            elif summary_type == 'extractive':
                summary = bert_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
            else:
                summary = "Invalid summary type selected."

            summary_filename = "webpage_summary.txt"
            
            with open(summary_filename, 'w', encoding="utf-8") as summary_file:
                summary_file.write(summary)

            return jsonify({
                "success": True, 
                "message": f"Summary generated successfully", 
                "summary": summary
            })
        except Exception as e:
            return jsonify({"success": False, "message": f"Error during summarization: {str(e)}"}), 400

    except Exception as e:
        return jsonify({"success": False, "message": f"Server error: {str(e)}"}), 500



@app.route("/video-upload", methods=["POST"])
def video_upload():
    try:
        data = request.get_json()
        if not data:
            return jsonify({"success": False, "message": "No data received"}), 400

        video_url = data.get("link")
        if not video_url:
            return jsonify({"success": False, "message": "Video URL is required"}), 400

        summary_type = data.get("summary_type", 'abstractive')
        summary_length = data.get("summary_length", 'medium')

        # Extract video ID from URL
        try:
            if "youtube.com" in video_url:
                video_id = video_url.split('v=')[-1].split('&')[0]
            elif "youtu.be" in video_url:
                video_id = video_url.split('/')[-1].split('?')[0]
            else:
                return jsonify({"success": False, "message": "Invalid YouTube URL format"}), 400
        except Exception:
            return jsonify({"success": False, "message": "Could not extract video ID from URL"}), 400

        try:
            transcript = YouTubeTranscriptApi.get_transcript(video_id)
            formatter = TextFormatter()
            formatted_transcript = formatter.format_transcript(transcript)
        except Exception as e:
            return jsonify({"success": False, "message": f"Error fetching transcript: {str(e)}"}), 400
        
        transcript_filename = f"{video_id}_transcript.txt"
        with open(transcript_filename, 'w', encoding="utf-8") as transcript_file:
            transcript_file.write(formatted_transcript)

        # Process the transcript to summarize
        max_input_length = 512

        if summary_length == 'short':
            max_output_length = 50
        elif summary_length == 'medium':
            max_output_length = 130
        elif summary_length == 'long':
            max_output_length = 250
        else:
            max_output_length = 130

        truncated_text = formatted_transcript[:max_input_length]

        try:
            if summary_type == 'abstractive':
                summary = bart_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
            elif summary_type == 'extractive':
                summary = bert_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
            else:
                summary = "Invalid summary type selected."

            summary_filename = f"{transcript_filename}_summary.txt"
            
            with open(summary_filename, 'w', encoding="utf-8") as summary_file:
                summary_file.write(summary)
            try:
                yt_info = YoutubeDL().extract_info(video_url, download=False)
                title = yt_info.get("title", "")
                search = VideosSearch(title, limit=6)
                results = search.result()

                recommendations = []
                for video in results["result"]:
                    if video["id"] != video_id:
                        recommendations.append({
                            "title": video["title"],
                            "videoId": video["id"],
                            "thumbnail": video["thumbnails"][0]["url"],
                        })
                        if len(recommendations) == 5:
                            break
            except Exception as e:
                print("Error during recommendation:", e)
                recommendations = []
            return jsonify({
                "success": True,
                "message": "Summary generated successfully",
                "summary": summary,
                "recommendations": recommendations
            })
        except Exception as e:
            return jsonify({"success": False, "message": f"Error during summarization: {str(e)}"}), 400

    except Exception as e:
        return jsonify({"success": False, "message": f"Server error: {str(e)}"}), 500


@app.route("/txt-summarize", methods=["POST"])
def txt_summarize():
    try:
        text_file = request.files.get("text_file")
        filename = text_file.filename
        text_file.save(filename)
        summary_type = request.form.get("summary_type")
        summary_length = request.form.get("summary_length")

        if not text_file:
            return jsonify({"message": "No file uploaded."}), 400
        
        with open(filename, 'r') as file:
            text = file.read()

        max_input_length = 512

        if summary_length == 'short':
            max_output_length = 50
        elif summary_length == 'medium':
            max_output_length = 130
        elif summary_length == 'long':
            max_output_length = 250
        else:
            max_output_length = 130

        truncated_text = text[:max_input_length]

        try:
            if summary_type == 'abstractive':
                summary = bart_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
            elif summary_type == 'extractive':
                summary = bert_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
            else:
                summary = "Invalid summary type selected."

            summary_filename = f"{filename}_summary.txt"
            
            with open(summary_filename, 'w', encoding="utf-8") as summary_file:
                summary_file.write(summary)

            message = f"Summary generated and saved to {summary_filename} using {summary_type.upper()} with {summary_length} length."
            return jsonify(success=True, message=message, summary=summary)
        except Exception as e:
            return jsonify(success=False, message=f"Error during summarization: {e}")

    except Exception as e:
        return jsonify({"message": str(e)}), 500




@app.route('/summarize', methods=['POST'])
def summarize():
    print("\n\n\n\nHello\n\n\n\n")
    filename = request.form.get('filename')
    summary_type = request.form.get('summary_type')
    summary_length = request.form.get('summary_length')

    print("\n\n\nWagwan\n\n\n")

    with open(filename, 'r', encoding="utf-8") as file:
        text = file.read()

    max_input_length = 512  # Model's input token limit

    if summary_length == 'short':
        max_output_length = 50
    elif summary_length == 'medium':
        max_output_length = 130
    elif summary_length == 'long':
        max_output_length = 250
    else:
        max_output_length = 130

    truncated_text = text[:max_input_length]

    try:
        if summary_type == 'abstractive':
            summary = bart_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
        elif summary_type == 'extractive':
            summary = bert_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
        else:
            summary = "Invalid summary type selected."

        summary_filename = f"{filename}_summary.txt"
        with open(summary_filename, 'w', encoding="utf-8") as summary_file:
            summary_file.write(summary)

        message = f"Summary generated and saved to {summary_filename} using {summary_type.upper()} with {summary_length} length."
    except Exception as e:
        message = f"Error during summarization: {e}"

    return render_template("ack.html", message=message)



# @app.route('/quiz', methods=['GET', 'POST'])
# def quiz():
#     global questions
#     if request.method == 'POST':
#         try:
#             uploaded_file = request.files['quiz']
#             filename = uploaded_file.filename
#             uploaded_file.save(filename)

#             with open(filename, 'r') as file:
#                 text = file.read()

#             questions = txt2questions(text)
#             print(questions)
#             return jsonify({"success": True, "message": "Quiz generated successfully!"})

#         except Exception as e:
#             return jsonify({"success": False, "message": f"Error processing quiz: {e}"}), 400

#     elif request.method == 'GET':
#         if questions:
#             return jsonify({"success": True, "questions": questions})
#         else:
#             return jsonify({"success": False, "message": "No quiz found. Please generate a quiz first."}), 404

@app.route('/quiz', methods=['GET', 'POST'])
def quiz():
    global questions

    if request.method == 'POST':
        if 'file' not in request.files:
            return jsonify({'success': False, 'message': 'No file uploaded'}), 400

        file = request.files['file']
        if file.filename == '':
            return jsonify({'success': False, 'message': 'Empty filename'}), 400

        try:
            text = file.read().decode('utf-8')

            prompt = f"""
            Generate exactly 10 multiple-choice questions (MCQs) from the text below.
            Return the result in the following Python dictionary format:

            {{
                1: {{
                    "question": "...",
                    "answer": "...",
                    "options": ["...", "...", "...", "..."]
                }},
                2: {{ ... }},
                ...
            }}

            Text:
            {text}
            """

            # Generate MCQs using Gemini
            response = model.generate_content(prompt)
            content = response.text.strip()

            # Safely evaluate the dictionary (assuming Gemini follows prompt format)
            questions = eval(content)  # NOTE: safe only if you trust Gemini's output

            return jsonify({"success": True, "message": "Quiz generated successfully!", "questions": questions})

        except Exception as e:
            return jsonify({'success': False, 'message': f'Error processing quiz: {e}'}), 500

    elif request.method == 'GET':
        if questions:
            return jsonify({"success": True, "questions": questions})
        else:
            return jsonify({"success": False, "message": "No quiz found. Please generate a quiz first."}), 404




@app.route('/results', methods=['POST'])
def result():
    answers = request.json
    correct_q = 0
    for q_num, data in questions.items():
        user_answer = answers.get(str(q_num))
        if user_answer == data['answer']:
            correct_q += 1

    return jsonify({
        "success": True,
        "correct": correct_q,
        "total": len(questions),
    })


@app.route('/translate', methods=['POST'])
def translate():
    try:
        data = request.get_json()
        text = data.get('text')
        target_language = data.get('target_language', 'en')

        if not text:
            return jsonify({
                "success": False,
                "message": "No text provided for translation"
            }), 400

        translator = Translator()
        translated = translator.translate(text, dest=target_language)

        return jsonify({
            "success": True,
            "translated_text": translated.text
        })
    except Exception as e:
        return jsonify({
            "success": False,
            "message": f"Translation error: {str(e)}"
        }), 500

@app.route("/get_recommendations", methods=["POST"])
def get_recommendations():
    print("\n\nInside /get_recommendations endpoint\n\n")
    data = request.get_json()
    video_url = data.get("video", "")
    print("Received video_url:", video_url)

    # Extract video ID from URL
    match = re.search(r"(?:v=|youtu\.be/)([a-zA-Z0-9_-]{11})", video_url)
    if not match:
        return jsonify({"error": "Invalid YouTube URL"}), 400

    video_id = match.group(1)

    try:
        # Get original video title
        with YoutubeDL({'quiet': True}) as ydl:
            info = ydl.extract_info(video_url, download=False)
            title = info.get("title", "")
    except Exception as e:
        print("yt-dlp failed:", str(e))
        return jsonify({"error": "Could not fetch video title."}), 500

    try:
        # Search for similar videos using yt-dlp
        search_url = f"ytsearch10:{title}"  # get top 10 to filter by popularity
        with YoutubeDL({'quiet': True}) as ydl:
            search_results = ydl.extract_info(search_url, download=False)['entries']

        recommendations = []

        for video in search_results:
            if video['id'] != video_id:  # Skip the original video
                views = video.get('view_count', 0)
                likes = video.get('like_count', 0)

                recommendations.append({
                    "title": video.get('title'),
                    "videoId": video.get('id'),
                    "thumbnail": video.get('thumbnail'),
                    "length": video.get('duration'),
                    "channel": video.get('uploader'),
                    "views": views,
                    "likes": likes,
                    "popularity": (views or 0) + (likes or 0)  # Simple popularity score
                })
                if len(recommendations) == 20:
                    break

        # Sort videos by popularity
        recommendations.sort(key=lambda x: x["popularity"], reverse=True)

        # Keep top 6 only
        final_recommendations = recommendations[:5]

        print("Recommended videos:", final_recommendations)
        return jsonify({"success": True, "recommendedVideos": final_recommendations})

    except Exception as e:
        print("Error during recommendations:", str(e))
        return jsonify({"error": str(e)}), 500
    
def extract_video_id(url):
    try:
        parsed_url = urllib.parse.urlparse(url)
        if parsed_url.hostname == 'youtu.be':
            return parsed_url.path[1:]
        if parsed_url.hostname in ['www.youtube.com', 'youtube.com']:
            query = urllib.parse.parse_qs(parsed_url.query)
            return query.get('v', [None])[0]
    except Exception:
        return None

def get_transcript(video_id, segment_count=4):
    transcript = YouTubeTranscriptApi.get_transcript(video_id)
    full_text = " ".join([entry['text'] for entry in transcript])
    words = full_text.split()
    chunk_size = len(words) // segment_count
    return [" ".join(words[i*chunk_size:(i+1)*chunk_size]) for i in range(segment_count)]


def compute_cosine_similarity(title, text_parts):
    similarities = []
    for part in text_parts:
        vectorizer = TfidfVectorizer().fit_transform([title, part])
        sim = cosine_similarity(vectorizer[0:1], vectorizer[1:2])[0][0]
        similarities.append(sim)
    return np.mean(similarities)

def analyze_comments(comments):
    polarity_counts = {'positive': 0, 'negative': 0, 'clickbait_phrases': 0}
    clickbait_keywords = [
        'fake', 'bullshit', 'hoax', 'wrong', 'liar', 'false', 'misinformation',
        'rumor', 'clickbait', 'myth', 'not true'
    ]
    for c in comments:
        text = c.lower()
        blob = TextBlob(text)
        if blob.sentiment.polarity > 0:
            polarity_counts['positive'] += 1
        elif blob.sentiment.polarity < 0:
            polarity_counts['negative'] += 1
        if any(kw in text for kw in clickbait_keywords):
            polarity_counts['clickbait_phrases'] += 1
    return polarity_counts, len(comments)

def get_video_details(video_id):
    youtube = build('youtube', 'v3', developerKey=YOUTUBE_API_KEY)
    request = youtube.videos().list(part="snippet,statistics", id=video_id)
    response = request.execute()
    data = response['items'][0]
    snippet = data['snippet']
    stats = data['statistics']
    return {
        'title': snippet['title'],
        'channel_id': snippet['channelId'],
        'likes': int(stats.get('likeCount', 0)),
        'dislikes': int(stats.get('dislikeCount', 1)),
        'views': int(stats.get('viewCount', 0)),
        'comment_count': int(stats.get('commentCount', 0))
    }

def get_channel_details(channel_id):
    youtube = build('youtube', 'v3', developerKey=YOUTUBE_API_KEY)
    request = youtube.channels().list(part="snippet,statistics", id=channel_id)
    response = request.execute()
    data = response['items'][0]
    stats = data['statistics']
    return {
        'subscribers': int(stats.get('subscriberCount', 1)),
        'views': int(stats.get('viewCount', 0)),
        'videos': int(stats.get('videoCount', 1))
    }

def get_top_comments(video_id, max_comments=50):
    youtube = build('youtube', 'v3', developerKey=YOUTUBE_API_KEY)
    request = youtube.commentThreads().list(
        part="snippet", videoId=video_id, maxResults=max_comments, textFormat="plainText"
    )
    response = request.execute()
    return [item['snippet']['topLevelComment']['snippet']['textDisplay'] for item in response['items']]

@app.route('/detect_clickbait', methods=['POST'])
def detect_clickbait():
    data = request.get_json()
    video_url = data.get('video_url')
    print(f"\nReceived video_url: {video_url}")
    if not video_url:
        return jsonify({'error': 'Missing video_url parameter'}), 400

    video_id = extract_video_id(video_url)
    if not video_id:
        return jsonify({'error': 'Invalid YouTube URL'}), 400
    print(f"\nVideo URL: {video_url}")
    print(f"\nVideo ID: {video_id}") 
    
    try:
        # Get data
        api_url = f"https://returnyoutubedislikeapi.com/votes?videoId={video_id}"
        response = requests.get(api_url)
        data = response.json()
        dislikes=data.get("dislikes")
        video_data = get_video_details(video_id)
        channel_data = get_channel_details(video_data['channel_id'])
        transcript_parts = get_transcript(video_id)
        avg_cs = compute_cosine_similarity(video_data['title'], transcript_parts)
        print(dislikes)
        print(video_data['likes'])
        dl_ratio = dislikes/ (video_data['likes'] + 1)

        # Comment analysis
        comments = get_top_comments(video_id)
        comment_stats, total_comments = analyze_comments(comments)
        fccr = comment_stats['clickbait_phrases'] / total_comments if total_comments else 0

        # Evidence checks
        evidence_1 = avg_cs > 0.10 or dl_ratio >= 0.40
        evidence_2 = fccr > 0.10
        evidence_3 = False  # Optional enhancement

        label = "Clickbait" if (evidence_1 or evidence_2 or evidence_3) else "Non-Clickbait"

        return jsonify({
            'label': label,
            'title': video_data['title'],
            'avg_cosine_similarity': round(avg_cs, 3),
            'dislike_like_ratio': round(dl_ratio, 3),
            'fake_comment_ratio': round(fccr, 3)
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500


if __name__ == '__main__':
    app.run(debug=True, port=5000)
